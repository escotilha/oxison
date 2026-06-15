"""Presentation adapter — slide text + speaker notes via python-pptx."""
from __future__ import annotations

from pathlib import Path

from .base import AdapterAvailability, SourceResult, SourceUnit


class PptxAdapter:
    name = "pptx"

    def detect(self, path: Path) -> bool:
        return path.suffix.lower() == ".pptx"

    def available(self) -> AdapterAvailability:
        try:
            import pptx  # noqa: F401
        except ImportError:
            return AdapterAvailability(
                available=False, reason="python-pptx not installed (pip install 'oxi-son[pptx]')"
            )
        return AdapterAvailability(available=True)

    def extract(self, path: Path) -> SourceResult:
        avail = self.available()
        if not avail.available:
            return SourceResult.skip(self.name, str(path), reason=avail.reason or "unavailable")
        from pptx import Presentation

        prs = Presentation(str(path))
        units: list[SourceUnit] = []
        for i, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
            if slide.has_notes_slide:
                note = slide.notes_slide.notes_text_frame.text
                if note.strip():
                    parts.append(f"[speaker notes] {note}")
            text = "\n".join(parts)
            if text.strip():
                units.append(
                    SourceUnit(
                        text=text,
                        source_type=self.name,
                        origin_path=str(path),
                        locator=f"pptx:{path.name}#slide-{i}",
                        metadata={"slide": i},
                    )
                )
        if not units:
            return SourceResult.skip(self.name, str(path), reason="no slide text")
        return SourceResult.ok(self.name, str(path), units=units)
