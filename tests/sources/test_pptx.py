from pathlib import Path

import pytest

from oxison.sources.pptx import PptxAdapter

pptx_mod = pytest.importorskip("pptx")


def _make_deck(path: Path) -> None:
    from pptx import Presentation
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = "Vision"
    s1.notes_slide.notes_text_frame.text = "speaker note one"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Roadmap"
    prs.save(str(path))


def test_pptx_detect(tmp_path: Path):
    assert PptxAdapter().detect(tmp_path / "deck.pptx")
    assert not PptxAdapter().detect(tmp_path / "x.pdf")


def test_pptx_extracts_slides_and_notes(tmp_path: Path):
    f = tmp_path / "deck.pptx"
    _make_deck(f)
    res = PptxAdapter().extract(f)
    assert res.status == "ok"
    assert res.unit_count == 2
    assert res.units[0].locator == "pptx:deck.pptx#slide-1"
    assert "Vision" in res.units[0].text
    assert "speaker note one" in res.units[0].text
    assert "Roadmap" in res.units[1].text


def test_pptx_skips_empty_deck(tmp_path: Path):
    f = tmp_path / "blank.pptx"
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # layout 6 = blank, no title placeholder
    prs.save(str(f))
    res = PptxAdapter().extract(f)
    assert res.status == "skipped"
    assert res.reason == "no slide text"
    assert res.unit_count == 0


def test_pptx_degrades_when_lib_absent(tmp_path: Path, monkeypatch):
    import builtins
    real_import = builtins.__import__

    def fake_import(name, *args, **kwargs):
        if name == "pptx":
            raise ImportError("forced: python-pptx absent")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", fake_import)
    f = tmp_path / "deck.pptx"
    f.write_bytes(b"PK fake")
    res = PptxAdapter().extract(f)
    assert res.status == "skipped"
    assert "python-pptx not installed" in (res.reason or "")
