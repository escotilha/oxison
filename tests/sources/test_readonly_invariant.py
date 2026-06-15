"""oxison's #1 invariant: source adapters READ inputs, never mutate them.

Each adapter's input file is hashed before and after extract(); the hash
must be unchanged. The AI-worker read-only guarantee is tested elsewhere
(the --allowedTools assertion); this covers the extraction adapters.
"""
from __future__ import annotations

import hashlib
from pathlib import Path

import pytest

from oxison.sources.docs import DocsAdapter


def _sha(p: Path) -> str:
    return hashlib.sha256(p.read_bytes()).hexdigest()


def test_docs_adapter_does_not_mutate_input(tmp_path: Path):
    f = tmp_path / "n.md"
    f.write_text("hello", encoding="utf-8")
    before = _sha(f)
    DocsAdapter().extract(f)
    assert _sha(f) == before


def test_pdf_adapter_does_not_mutate_input(tmp_path: Path):
    pytest.importorskip("pypdf")
    from pypdf import PdfWriter
    f = tmp_path / "x.pdf"
    w = PdfWriter()
    w.add_blank_page(width=100, height=100)
    with open(f, "wb") as fh:
        w.write(fh)
    before = _sha(f)
    from oxison.sources.pdf import PdfAdapter
    PdfAdapter().extract(f)
    assert _sha(f) == before


def test_pptx_adapter_does_not_mutate_input(tmp_path: Path):
    pytest.importorskip("pptx")
    from pptx import Presentation
    f = tmp_path / "d.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "t"
    prs.save(str(f))
    before = _sha(f)
    from oxison.sources.pptx import PptxAdapter
    PptxAdapter().extract(f)
    assert _sha(f) == before


def test_docx_adapter_does_not_mutate_input(tmp_path: Path):
    pytest.importorskip("docx")
    from docx import Document
    f = tmp_path / "s.docx"
    d = Document()
    d.add_paragraph("requirement one")
    d.save(str(f))
    before = _sha(f)
    from oxison.sources.docx import DocxAdapter
    DocxAdapter().extract(f)
    assert _sha(f) == before
